#!/usr/bin/env python3
"""Read a corporate .pptx template and emit a starter brand.json.

This is what turns a brand into configuration. Point it at your organisation's
corporate template and it reports what is actually in there: every slide master,
every layout name, the slide size, the theme palette, the fonts, and which
layouts carry a full-bleed background picture (the usual way a "wave" or
decorative background is delivered).

    python3 make_brand_pack.py corporate.pptx --out brand/brand.json

The generated file is a starting point, not an answer. The alignment grid and
the reserved logo zone are measured from one layout you nominate with
--grid-from; check them and correct them by eye. Everything the script cannot
determine is written as null so it is obvious what still needs a human.

Requires python-pptx.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:  # pragma: no cover
    print("python-pptx is required:  pip install python-pptx", file=sys.stderr)
    raise SystemExit(2)


def emu_to_in(v) -> float | None:
    return None if v is None else round(Emu(int(v)).inches, 3)


def layout_has_full_bleed_picture(layout, slide_w: int, slide_h: int) -> bool:
    """True when a layout carries a picture covering ~the whole slide.

    That is how decorative backgrounds are almost always shipped. It matters
    because a slide using such a layout must NOT have its background overridden:
    setting an explicit fill hides the very decoration you selected the layout
    for.
    """
    for shape in layout.shapes:
        if shape.shape_type is None or "PICTURE" not in str(shape.shape_type):
            continue
        if shape.width is None or shape.height is None:
            continue
        if shape.width >= slide_w * 0.95 and shape.height >= slide_h * 0.95:
            return True
    return False


def collect_masters(prs) -> list[dict]:
    """Every master and its layouts.

    prs.slide_layouts only exposes the FIRST master. Templates routinely put
    dark-mode or alternate layouts on a second master, so anything that walks
    only prs.slide_layouts will raise "layout not found" on half the template.
    """
    sw, sh = prs.slide_width, prs.slide_height
    out = []
    for mi, master in enumerate(prs.slide_masters):
        layouts = []
        for layout in master.slide_layouts:
            layouts.append({
                "name": layout.name,
                "placeholders": sorted(
                    {str(ph.placeholder_format.type) for ph in layout.placeholders}
                ),
                "full_bleed_background_picture":
                    layout_has_full_bleed_picture(layout, sw, sh),
            })
        out.append({"index": mi, "name": master.name, "layouts": layouts})
    return out


def theme_palette(prs) -> dict:
    """Pull the theme colour scheme straight out of the master's XML."""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    palette: dict[str, str] = {}
    try:
        theme = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        root = theme._element if hasattr(theme, "_element") else None
        if root is None:
            from lxml import etree
            root = etree.fromstring(theme.blob)
        for el in root.findall(".//a:clrScheme/*", ns):
            tag = el.tag.split("}")[-1]
            srgb = el.find("a:srgbClr", ns)
            sysc = el.find("a:sysClr", ns)
            if srgb is not None:
                palette[tag] = "#" + srgb.get("val")
            elif sysc is not None and sysc.get("lastClr"):
                palette[tag] = "#" + sysc.get("lastClr")
    except Exception as exc:  # noqa: BLE001 - a missing theme is not fatal
        palette["_error"] = f"could not read theme: {exc}"
    return palette


def theme_fonts(prs) -> dict:
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    fonts: dict[str, str | None] = {"major": None, "minor": None}
    try:
        theme = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        root = theme._element if hasattr(theme, "_element") else None
        if root is None:
            from lxml import etree
            root = etree.fromstring(theme.blob)
        for key, path in (("major", "a:majorFont/a:latin"), ("minor", "a:minorFont/a:latin")):
            el = root.find(f".//a:fontScheme/{path}", ns)
            if el is not None:
                fonts[key] = el.get("typeface")
    except Exception:  # noqa: BLE001
        pass
    return fonts


def grid_from_layout(prs, layout_name: str) -> dict:
    """Measure an alignment grid from one nominated layout.

    Placeholder bounds on a real content layout are the most reliable source of
    the left margin, title band, and content top. Treat the output as a first
    read to be checked by eye, not as truth.
    """
    target = None
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                target = layout
                break
    if target is None:
        return {"_error": f"layout not found: {layout_name!r}"}

    grid = {
        "slide_w_in": emu_to_in(prs.slide_width),
        "slide_h_in": emu_to_in(prs.slide_height),
        "measured_from_layout": layout_name,
    }
    lefts = []
    for ph in target.placeholders:
        kind = str(ph.placeholder_format.type)
        entry = {
            "x_in": emu_to_in(ph.left), "y_in": emu_to_in(ph.top),
            "w_in": emu_to_in(ph.width), "h_in": emu_to_in(ph.height),
        }
        if ph.left is not None:
            lefts.append(ph.left)
        if "TITLE" in kind and "title" not in grid:
            grid["title"] = entry
        elif "BODY" in kind and "body" not in grid:
            grid["body"] = entry
    grid["left_edge_in"] = emu_to_in(min(lefts)) if lefts else None
    return grid


def reserved_zones(prs, master_index: int = 0) -> list[dict]:
    """Picture shapes on a master, the usual home of a logo.

    Anything listed here is space your content must not cover.
    """
    zones = []
    try:
        master = prs.slide_masters[master_index]
    except IndexError:
        return zones
    for shape in master.shapes:
        if shape.shape_type is None or "PICTURE" not in str(shape.shape_type):
            continue
        zones.append({
            "name": shape.name,
            "x_in": emu_to_in(shape.left), "y_in": emu_to_in(shape.top),
            "w_in": emu_to_in(shape.width), "h_in": emu_to_in(shape.height),
            "note": "Confirm this is the logo before treating it as reserved.",
        })
    return zones


def build(template: Path, grid_layout: str | None) -> dict:
    prs = Presentation(str(template))
    masters = collect_masters(prs)
    decorated = [
        {"master": m["index"], "layout": l["name"]}
        for m in masters for l in m["layouts"]
        if l["full_bleed_background_picture"]
    ]
    return {
        "$comment": (
            "Starter brand pack. Values under 'grid' and 'reserved_zones' are measured "
            "and must be checked by eye. Nulls mark what still needs a human."
        ),
        "template": {
            "build_base": str(template),
            "slide_w_in": emu_to_in(prs.slide_width),
            "slide_h_in": emu_to_in(prs.slide_height),
            "master_count": len(masters),
            "layout_count": sum(len(m["layouts"]) for m in masters),
        },
        "typography": {
            "theme_fonts": theme_fonts(prs),
            "min_body_pt": 14,
        },
        "palette": theme_palette(prs),
        "grid": grid_from_layout(prs, grid_layout) if grid_layout else None,
        "reserved_zones": reserved_zones(prs),
        "modes": {
            "light": {"layout_hints": [], "set_background": False},
            "dark_flat": {"layout_hints": [], "background_hex": None, "set_background": True},
            "dark_decorated": {
                "layout_hints": [d["layout"] for d in decorated],
                "set_background": False,
                "note": "These layouts carry their own background picture. Never set an "
                        "explicit slide background on them or the decoration disappears.",
            },
        },
        "masters": masters,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("template", type=Path, help="Corporate .pptx template to read.")
    ap.add_argument("--out", type=Path, help="Write JSON here instead of stdout.")
    ap.add_argument("--grid-from", metavar="LAYOUT_NAME",
                    help="Measure the alignment grid from this layout. "
                         "Pick a plain one-column content layout.")
    ap.add_argument("--list-layouts", action="store_true",
                    help="Print layout names and exit, to help choose --grid-from.")
    args = ap.parse_args()

    if not args.template.exists():
        print(f"No such file: {args.template}", file=sys.stderr)
        return 2

    if args.list_layouts:
        prs = Presentation(str(args.template))
        for mi, master in enumerate(prs.slide_masters):
            print(f"Master {mi}: {master.name}")
            for layout in master.slide_layouts:
                print(f"    {layout.name}")
        return 0

    pack = build(args.template, args.grid_from)
    text = json.dumps(pack, indent=2)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(text + "\n")
        t = pack["template"]
        print(f"Wrote {args.out}")
        print(f"  {t['master_count']} master(s), {t['layout_count']} layouts, "
              f"{t['slide_w_in']}x{t['slide_h_in']} in")
        print(f"  {len(pack['modes']['dark_decorated']['layout_hints'])} layout(s) carry a "
              f"full-bleed background picture")
        print(f"  {len(pack['reserved_zones'])} candidate reserved zone(s) to confirm")
        if pack["grid"] is None:
            print("  grid: not measured (pass --grid-from to measure one)")
    else:
        print(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
