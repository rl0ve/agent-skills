#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def delete_slide(prs: Presentation, idx: int) -> None:
    slide_id = prs.slides._sldIdLst[idx]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    del prs.slides._sldIdLst[idx]


def find_layout(prs: Presentation, layout_name: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == layout_name:
                return layout
    return None


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Start a new template-backed deck from a build-base file and add "
            "slides using named existing layouts."
        )
    )
    parser.add_argument("build_base", help="Path to the brand build-base .pptx")
    parser.add_argument("output", help="Path to write the new .pptx")
    parser.add_argument(
        "--layout",
        dest="layouts",
        action="append",
        required=True,
        help="Exact slide layout name to add. Repeat for multiple slides.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    build_base = Path(args.build_base).expanduser().resolve()
    output = Path(args.output).expanduser().resolve()

    if not build_base.exists():
        raise SystemExit(f"Build-base template not found: {build_base}")

    prs = Presentation(str(build_base))

    missing = [name for name in args.layouts if find_layout(prs, name) is None]
    if missing:
        available = sorted(
            {
                layout.name
                for master in prs.slide_masters
                for layout in master.slide_layouts
            }
        )
        available_list = "\n".join(f"- {name}" for name in available)
        missing_list = ", ".join(missing)
        raise SystemExit(
            f"Requested layout(s) not found: {missing_list}\nAvailable layouts:\n{available_list}"
        )

    for idx in range(len(prs.slides) - 1, -1, -1):
        delete_slide(prs, idx)

    for layout_name in args.layouts:
        prs.slides.add_slide(find_layout(prs, layout_name))

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
