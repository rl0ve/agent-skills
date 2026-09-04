#!/usr/bin/env python3
"""Check a .pptx for the things that actually go wrong in a generated deck.

Every check here exists because it catches a defect you cannot see by reading
the build script, only by looking at the deck. Running this is not a substitute
for looking at rendered slides; it is what makes looking cheap, by telling you
which slides to look at first.

    python3 check_deck.py deck.pptx
    python3 check_deck.py deck.pptx --brand brand/brand.json
    python3 check_deck.py deck.pptx --json          # machine-readable
    python3 check_deck.py deck.pptx --only FIT,GEO  # one family
    python3 check_deck.py deck.pptx --min-pt 12 --max-same-layout 3

Exit codes: 0 clean or notes only, 1 at least one WARN, 2 at least one FAIL,
3 the file could not be read.

Checks
------
PKG001  duplicate package entries, the PowerPoint "repair" dialog        FAIL
PLC001  empty placeholders inherited from the layout (editing view only) NOTE
PLC002  template helper text left in a delivered slide                   FAIL
FIT001  text likely overflows its shape (estimated)                      WARN
FIT003  shape grows to fit its text, so it will not stay where drawn      WARN
FIT002  text below the minimum readable point size                       WARN
GEO001  a shape extends past the slide edge                              WARN
GEO002  an opaque shape covers text that sits below it                   FAIL
GEO003  a shape intrudes on a reserved zone, usually the logo            FAIL
GEO004  sibling shapes that nearly share a top edge but do not           NOTE
VAR001  one layout used more than the allowed number of times            WARN
VAR002  fewer than four distinct layouts across content slides           WARN

Requires python-pptx.
"""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
import zipfile
from collections import Counter
from dataclasses import dataclass, asdict, field
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:  # pragma: no cover
    print("python-pptx is required:  pip install python-pptx", file=sys.stderr)
    raise SystemExit(3)


FAIL, WARN, NOTE = "FAIL", "WARN", "NOTE"

# Text left behind by a template. Matching any of these in a delivered deck
# means a placeholder was styled but never actually filled in.
HELPER_TEXT = re.compile(
    r"lorem ipsum|click to (edit|add)|your (title|text|logo) here|"
    r"\[?(insert|add|replace|tbd|todo|xxx|placeholder)\b",
    re.I,
)

# Rough width of one character as a fraction of the font's point size, for a
# proportional sans face. Used only to estimate line wrapping.
CHAR_WIDTH_RATIO = 0.50
LINE_HEIGHT_RATIO = 1.22
DEFAULT_PT = 18.0


@dataclass
class Finding:
    code: str
    severity: str
    slide: int
    message: str
    shape: str | None = None

    def line(self) -> str:
        where = f"slide {self.slide}" if self.slide else "deck"
        what = f" [{self.shape}]" if self.shape else ""
        return f"{self.severity}  {self.code}  {where}{what}: {self.message}"


@dataclass
class Brand:
    reserved: list[dict] = field(default_factory=list)
    min_pt: float | None = None

    @classmethod
    def load(cls, path: Path | None) -> "Brand":
        if not path:
            return cls()
        data = json.loads(path.read_text())
        return cls(
            reserved=data.get("reserved_zones") or [],
            min_pt=(data.get("typography") or {}).get("min_body_pt"),
        )


def emu_in(v) -> float:
    return Emu(int(v)).inches if v is not None else 0.0


def walk(shapes):
    """Yield every shape, descending into groups, preserving z-order."""
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) is not None and str(shape.shape_type) == "GROUP (6)":
            try:
                yield from walk(shape.shapes)
            except (AttributeError, NotImplementedError):
                pass


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return " ".join((shape.text_frame.text or "").split())


def effective_pt(paragraph, shape) -> float:
    """Best available font size for a paragraph, in points.

    python-pptx returns None wherever the size is inherited from the layout or
    master, which is most of the time on a template-backed deck. Walk what is
    knowable and fall back to a stated default rather than guessing silently.
    """
    for run in paragraph.runs:
        if run.font.size is not None:
            return run.font.size.pt
    if paragraph.font.size is not None:
        return paragraph.font.size.pt
    try:
        ph = shape.text_frame.paragraphs[0]
        if ph.font.size is not None:
            return ph.font.size.pt
    except (AttributeError, IndexError):
        pass
    return DEFAULT_PT


def estimate_text_height_in(shape) -> tuple[float, bool]:
    """Estimate rendered text height in inches.

    Returns (height, exact). `exact` is always False: this wraps text by
    character count against an average glyph width, which is an approximation.
    It is reliable for catching text that badly overruns its box and unreliable
    for anything within about 15% of the boundary, so callers should treat a
    near miss as "go and look", never as a hard failure.
    """
    tf = shape.text_frame
    width_in = emu_in(shape.width)
    try:
        width_in -= emu_in(tf.margin_left) + emu_in(tf.margin_right)
    except (AttributeError, TypeError):
        pass
    if width_in <= 0:
        return 0.0, False

    total_in = 0.0
    for para in tf.paragraphs:
        text = "".join(r.text for r in para.runs) or ""
        pt = effective_pt(para, shape)
        char_w_in = (pt * CHAR_WIDTH_RATIO) / 72.0
        per_line = max(1, int(width_in / char_w_in)) if char_w_in > 0 else 1
        lines = max(1, math.ceil(len(text) / per_line)) if text else 1
        total_in += lines * (pt * LINE_HEIGHT_RATIO) / 72.0

    try:
        total_in += emu_in(tf.margin_top) + emu_in(tf.margin_bottom)
    except (AttributeError, TypeError):
        pass
    return total_in, False


def autofit_mode(shape) -> str:
    """How PowerPoint reconciles text and box: 'none', 'shrink-text', or 'grow-shape'.

    These fail in different ways and must not be collapsed into one boolean.
    'shrink-text' genuinely handles overflow, so there is nothing to report.
    'grow-shape' does not: the box silently becomes taller than the author drew
    it, which is how content ends up off the slide or on top of its neighbour.
    """
    try:
        mode = str(shape.text_frame.auto_size)
    except (AttributeError, ValueError):
        return "none"
    if "TEXT_TO_FIT_SHAPE" in mode:
        return "shrink-text"
    if "SHAPE_TO_FIT_TEXT" in mode:
        return "grow-shape"
    return "none"


def is_opaque(shape) -> bool:
    try:
        t = shape.fill.type
    except (AttributeError, TypeError, ValueError):
        return False
    if t is None:
        return False
    name = str(t)
    return "SOLID" in name or "PICTURE" in name or "GRADIENT" in name


def text_extent_rect(shape) -> tuple[float, float, float, float] | None:
    """The rectangle the shape's ink plausibly occupies, not its box.

    A wide, left-aligned text box whose text stops well short of the right edge
    does not actually collide with anything over there. Testing box bounds
    instead of text extent over-reports badly enough that people stop reading
    the findings, so estimate how far the longest line really reaches.
    Non-text shapes fall back to their full bounds, which for them is correct.
    """
    r = rect(shape)
    if r is None or not shape.has_text_frame:
        return r
    text = shape_text(shape)
    if not text:
        return None
    tf = shape.text_frame
    widest = 0.0
    for para in tf.paragraphs:
        line = "".join(run.text for run in para.runs)
        pt = effective_pt(para, shape)
        widest = max(widest, len(line) * (pt * CHAR_WIDTH_RATIO) / 72.0)
    try:
        align = str(tf.paragraphs[0].alignment or "")
    except (AttributeError, IndexError):
        align = ""
    box_w = r[2] - r[0]
    ink_w = min(box_w, widest)
    if "CENTER" in align:
        pad = (box_w - ink_w) / 2
        return (r[0] + pad, r[1], r[2] - pad, r[3])
    if "RIGHT" in align:
        return (r[2] - ink_w, r[1], r[2], r[3])
    return (r[0], r[1], r[0] + ink_w, r[3])


def rect(shape) -> tuple[float, float, float, float] | None:
    if None in (shape.left, shape.top, shape.width, shape.height):
        return None
    x, y = emu_in(shape.left), emu_in(shape.top)
    return (x, y, x + emu_in(shape.width), y + emu_in(shape.height))


def overlap_area(a, b) -> float:
    ox = max(0.0, min(a[2], b[2]) - max(a[0], b[0]))
    oy = max(0.0, min(a[3], b[3]) - max(a[1], b[1]))
    return ox * oy


def area(r) -> float:
    return max(0.0, r[2] - r[0]) * max(0.0, r[3] - r[1])


# --------------------------------------------------------------------------- #
# checks
# --------------------------------------------------------------------------- #

def check_package(path: Path) -> list[Finding]:
    counts: Counter = Counter()
    with zipfile.ZipFile(path) as z:
        for info in z.infolist():
            counts[info.filename] += 1
    dupes = [n for n, c in counts.items() if c > 1]
    if not dupes:
        return []
    return [Finding(
        "PKG001", FAIL, 0,
        f"{len(dupes)} duplicated package entr"
        f"{'y' if len(dupes) == 1 else 'ies'} ({', '.join(sorted(dupes)[:3])}"
        f"{'…' if len(dupes) > 3 else ''}). PowerPoint will show the repair dialog on "
        f"open. Fix with dedupe_pptx_zip.py.",
    )]


def check_slide(idx: int, slide, slide_w: float, slide_h: float,
                brand: Brand, min_pt: float) -> list[Finding]:
    out: list[Finding] = []
    shapes = list(walk(slide.shapes))
    empty_placeholders: list[str] = []

    for shape in shapes:
        name = shape.name
        text = shape_text(shape)
        r = rect(shape)

        # --- placeholders and leftover template text ---
        if shape.is_placeholder and not text:
            try:
                kind = str(shape.placeholder_format.type)
            except (AttributeError, ValueError):
                kind = "unknown"
            if kind not in ("SLIDE_NUMBER (13)", "FOOTER (15)", "DATE (16)"):
                empty_placeholders.append(f"{name} ({kind})")
        if text and HELPER_TEXT.search(text):
            out.append(Finding("PLC002", FAIL, idx,
                               f"template helper text still present: {text[:60]!r}", name))

        # --- typography ---
        if shape.has_text_frame and text:
            for para in shape.text_frame.paragraphs:
                pt = effective_pt(para, shape)
                if pt < min_pt:
                    out.append(Finding("FIT002", WARN, idx,
                                       f"text at {pt:g}pt is below the {min_pt:g}pt floor.", name))
                    break

            mode = autofit_mode(shape)
            if r and mode != "shrink-text":
                est, _ = estimate_text_height_in(shape)
                box_h = r[3] - r[1]
                if box_h > 0 and est > box_h * 1.15:
                    if mode == "none":
                        out.append(Finding("FIT001", WARN, idx,
                                           f"text is estimated at {est:.2f}in in a {box_h:.2f}in "
                                           f"box and will probably overflow. Estimated, so look at "
                                           f"the rendered slide before changing anything.", name))
                    else:  # grow-shape
                        grown_bottom = r[1] + est
                        off = grown_bottom > slide_h + 0.01
                        out.append(Finding("FIT003", WARN, idx,
                                           f"set to grow to fit its text, so the box drawn at "
                                           f"{box_h:.2f}in will render nearer {est:.2f}in"
                                           + (f" and run off the bottom of the slide."
                                              if off else
                                              f" and overlap whatever sits beneath it.")
                                           + " Size the text to the box, or the box to the text.",
                                           name))

        # --- geometry ---
        if r:
            if r[0] < -0.01 or r[1] < -0.01 or r[2] > slide_w + 0.01 or r[3] > slide_h + 0.01:
                out.append(Finding("GEO001", WARN, idx,
                                   f"extends past the slide edge "
                                   f"(x {r[0]:.2f}–{r[2]:.2f}in, y {r[1]:.2f}–{r[3]:.2f}in, "
                                   f"slide is {slide_w:.2f}x{slide_h:.2f}in).", name))

            # A template's own placeholders are allowed to occupy template space;
            # flagging them just trains people to ignore the check.
            ink = text_extent_rect(shape)
            for zone in ([] if shape.is_placeholder or ink is None else brand.reserved):
                try:
                    zr = (float(zone["x_in"]), float(zone["y_in"]),
                          float(zone["x_in"]) + float(zone["w_in"]),
                          float(zone["y_in"]) + float(zone["h_in"]))
                except (KeyError, TypeError, ValueError):
                    continue
                if overlap_area(ink, zr) > 0.0005:
                    out.append(Finding("GEO003", FAIL, idx,
                                       f"intrudes on reserved zone "
                                       f"{zone.get('name', 'unnamed')!r}. That space belongs to "
                                       f"the template, usually the logo.", name))

    # Empty placeholders do not render in presentation or print, only in the
    # editing view, so this is a tidiness note and not a defect. One finding per
    # slide: a template-canvas build inherits every layout placeholder, and one
    # finding each would bury everything else.
    if empty_placeholders:
        out.append(Finding("PLC001", NOTE, idx,
                           f"{len(empty_placeholders)} empty placeholder(s) inherited from the "
                           f"layout ({', '.join(empty_placeholders[:3])}"
                           f"{'…' if len(empty_placeholders) > 3 else ''}). They show as prompt "
                           f"text in the editing view. Remove the element, or use a layout with "
                           f"fewer placeholders."))

    # --- occlusion: an opaque shape drawn over text that sits below it ---
    for i, top in enumerate(shapes):
        if not is_opaque(top):
            continue
        tr = rect(top)
        if not tr or area(tr) <= 0:
            continue
        for below in shapes[:i]:
            btext = shape_text(below)
            if not btext:
                continue
            br = rect(below)
            if not br or area(br) <= 0:
                continue
            covered = overlap_area(tr, br) / area(br)
            if covered > 0.30:
                out.append(Finding("GEO002", FAIL, idx,
                                   f"{top.name!r} covers {covered * 100:.0f}% of the text shape "
                                   f"below it ({btext[:40]!r}). Move one of them.", below.name))
                break

    # --- near-miss alignment among top-level shapes ---
    tops = sorted({round(emu_in(s.top), 3) for s in slide.shapes if s.top is not None})
    for a, b in zip(tops, tops[1:]):
        if 0.005 < (b - a) < 0.06:
            out.append(Finding("GEO004", NOTE, idx,
                               f"two shape tops sit at {a:.3f}in and {b:.3f}in. That reads as "
                               f"a misalignment rather than a deliberate offset."))
            break

    return out


def check_variety(layout_names: list[str], max_same: int, min_distinct: int) -> list[Finding]:
    if len(layout_names) < 4:
        return []
    counts = Counter(layout_names)
    out: list[Finding] = []
    for name, n in counts.most_common():
        if n > max_same:
            out.append(Finding("VAR001", WARN, 0,
                               f"layout {name!r} is used on {n} of {len(layout_names)} content "
                               f"slides (limit {max_same}). A deck built from one pattern reads as "
                               f"repetitive however good the content is."))
    if len(counts) < min_distinct:
        out.append(Finding("VAR002", WARN, 0,
                           f"only {len(counts)} distinct layout(s) across {len(layout_names)} "
                           f"content slides. Use at least {min_distinct}."))
    return out


def check(path: Path, brand: Brand, min_pt: float,
          max_same: int, min_distinct: int) -> list[Finding]:
    # Package integrity first, and on its own. A package with duplicated entries
    # often will not parse at all, so opening it with python-pptx before this
    # check turns the one diagnosis that matters into an unrelated traceback.
    findings = check_package(path)
    if findings:
        findings.append(Finding(
            "PKG001", NOTE, 0,
            "Slide-level checks were skipped: fix the package first, then re-run. "
            "A duplicated package often cannot be parsed at all."))
        return findings

    prs = Presentation(str(path))
    slide_w, slide_h = emu_in(prs.slide_width), emu_in(prs.slide_height)

    layout_names: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        layout_names.append(slide.slide_layout.name)
        findings += check_slide(idx, slide, slide_w, slide_h, brand, min_pt)
    findings += check_variety(layout_names, max_same, min_distinct)
    return findings


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--brand", type=Path,
                    help="brand.json, for reserved zones and the point-size floor.")
    ap.add_argument("--min-pt", type=float, default=None,
                    help="Minimum readable point size (default 14, or the brand pack's value).")
    ap.add_argument("--max-same-layout", type=int, default=3,
                    help="How often one layout may repeat across content slides (default 3).")
    ap.add_argument("--min-distinct-layouts", type=int, default=4,
                    help="Distinct layouts required across content slides (default 4).")
    ap.add_argument("--only", help="Comma-separated check families to run, e.g. FIT,GEO,VAR,PLC,PKG.")
    ap.add_argument("--json", action="store_true", help="Emit findings as JSON.")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"No such file: {args.pptx}", file=sys.stderr)
        return 3

    brand = Brand.load(args.brand)
    min_pt = args.min_pt if args.min_pt is not None else (brand.min_pt or 14.0)

    try:
        findings = check(args.pptx, brand, min_pt,
                         args.max_same_layout, args.min_distinct_layouts)
    except Exception as exc:  # noqa: BLE001
        print(f"Could not read {args.pptx.name}: {exc}", file=sys.stderr)
        return 3

    if args.only:
        families = {f.strip().upper() for f in args.only.split(",") if f.strip()}
        findings = [f for f in findings if f.code[:3] in families]

    if args.json:
        print(json.dumps([asdict(f) for f in findings], indent=2))
    else:
        order = {FAIL: 0, WARN: 1, NOTE: 2}
        for f in sorted(findings, key=lambda f: (order[f.severity], f.slide, f.code)):
            print(f.line())
        counts = Counter(f.severity for f in findings)
        if not findings:
            print(f"{args.pptx.name}: clean.")
        else:
            print(f"\n{args.pptx.name}: "
                  f"{counts.get(FAIL, 0)} fail, {counts.get(WARN, 0)} warn, {counts.get(NOTE, 0)} note.")
            print("FIT001 is an estimate. Confirm it against a rendered slide before acting on it.")

    if any(f.severity == FAIL for f in findings):
        return 2
    if any(f.severity == WARN for f in findings):
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
