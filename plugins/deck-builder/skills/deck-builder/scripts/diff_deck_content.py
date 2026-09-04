#!/usr/bin/env python3
"""Report content present in an original deck and missing from a rebuilt one.

The number one failure when condensing or restructuring a deck: the build script
is rewritten for the new structure, and anything not explicitly coded into it
disappears without a word. Stats, proof points, customer examples, framing
language. Nobody notices until the person who supplied that content reads it.

    python3 diff_deck_content.py original.pptx rebuilt.pptx
    python3 diff_deck_content.py original.pptx rebuilt.pptx --json
    python3 diff_deck_content.py original.pptx rebuilt.pptx --all   # every phrase

By default this reports the content most expensive to lose:

  FACT   a number, percentage, money amount, or year that is gone
  NAME   a capitalised multi-word phrase that is gone, e.g. a product or customer
  PHRASE a whole sentence of eight or more words that is gone

Matching is whole-deck, not slide-by-slide, so moving content between slides is
not reported. Only content that left the deck entirely is.

Exit codes: 0 nothing lost, 1 something was dropped, 2 a file could not be read.

Requires python-pptx.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, asdict
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    print("python-pptx is required:  pip install python-pptx", file=sys.stderr)
    raise SystemExit(2)


# A number worth keeping: 70%, $1.2M, 5-8, 30+, 2026, 1,400
FACT = re.compile(r"(?<![\w.])(?:[$£€]\s?\d[\d,.]*\s?[KMB]?|\d[\d,.]*\s?%|\d[\d,.]*\s?[+-]\s?\d[\d,.]*|\d{4}|\d[\d,.]*)(?![\w.])")
# Two or more capitalised words in a row, the shape of a product or company name.
NAME = re.compile(r"\b(?:[A-Z][\w&.-]+)(?:\s+(?:[A-Z][\w&.-]+|of|and|for|the)){1,4}\b")

STOP_NAMES = {
    "The", "This", "That", "These", "Those", "We", "Our", "You", "Your", "It",
    "There", "What", "When", "Where", "Why", "How", "And", "But", "For", "With",
}


@dataclass
class Loss:
    kind: str
    value: str
    slides: list[int]


def normalise(text: str) -> str:
    return " ".join(text.replace("’", "'").replace("–", "-").split())


def deck_text(path: Path) -> dict[int, str]:
    prs = Presentation(str(path))
    out: dict[int, str] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        def collect(shapes):
            for shape in shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                parts.append(cell.text)
                if str(getattr(shape, "shape_type", "")) == "GROUP (6)":
                    try:
                        collect(shape.shapes)
                    except (AttributeError, NotImplementedError):
                        pass

        collect(slide.shapes)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
            parts.append(slide.notes_slide.notes_text_frame.text)
        # Separate shapes with a sentinel so a name or sentence match cannot run
        # across the join between two unrelated text boxes, which produced
        # phantom names like "Four Technical Differentiators One".
        out[idx] = normalise(" ¶ ".join(parts))
    return out


def sentences(text: str) -> list[str]:
    return [s.strip() for s in re.split(r"(?<=[.!?])\s+|\s*¶\s*", text)
            if len(s.split()) >= 8]


def extract(per_slide: dict[int, str]) -> dict[str, dict[str, list[int]]]:
    found: dict[str, dict[str, list[int]]] = {"FACT": {}, "NAME": {}, "PHRASE": {}}
    for idx, text in per_slide.items():
        for m in FACT.findall(text):
            found["FACT"].setdefault(m.strip(), []).append(idx)
        for chunk in text.split("¶"):
            for m in NAME.findall(chunk):
                m = m.strip()
                if not m or m.split()[0] in STOP_NAMES:
                    continue
                found["NAME"].setdefault(m, []).append(idx)
        for s in sentences(text):
            found["PHRASE"].setdefault(s, []).append(idx)
    return found


def compare(old: Path, new: Path, show_all: bool) -> list[Loss]:
    old_slides, new_slides = deck_text(old), deck_text(new)
    new_blob = " ".join(new_slides.values())
    new_blob_lower = new_blob.lower()
    old_found = extract(old_slides)

    losses: list[Loss] = []
    for kind in ("FACT", "NAME", "PHRASE"):
        for value, slides in old_found[kind].items():
            if kind == "PHRASE":
                # A sentence counts as surviving if most of its distinctive words
                # do; a rebuild legitimately rewords, and demanding an exact match
                # would report every surviving idea as lost.
                words = [w for w in re.findall(r"\w+", value.lower()) if len(w) > 3]
                if not words:
                    continue
                kept = sum(1 for w in set(words) if w in new_blob_lower)
                if kept / len(set(words)) >= 0.7:
                    continue
            elif value.lower() in new_blob_lower:
                continue
            losses.append(Loss(kind, value, sorted(set(slides))))

    if not show_all:
        # One representative phrase per original slide keeps the report readable.
        seen_phrase_slides: set[int] = set()
        trimmed: list[Loss] = []
        for loss in losses:
            if loss.kind == "PHRASE":
                key = loss.slides[0] if loss.slides else 0
                if key in seen_phrase_slides:
                    continue
                seen_phrase_slides.add(key)
            trimmed.append(loss)
        losses = trimmed

    order = {"FACT": 0, "NAME": 1, "PHRASE": 2}
    return sorted(losses, key=lambda l: (order[l.kind], l.slides, l.value))


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("original", type=Path)
    ap.add_argument("rebuilt", type=Path)
    ap.add_argument("--all", action="store_true",
                    help="Report every dropped phrase, not one per original slide.")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    for f in (args.original, args.rebuilt):
        if not f.exists():
            print(f"No such file: {f}", file=sys.stderr)
            return 2

    try:
        losses = compare(args.original, args.rebuilt, args.all)
    except Exception as exc:  # noqa: BLE001
        print(f"Could not compare: {exc}", file=sys.stderr)
        return 2

    if args.json:
        print(json.dumps([asdict(l) for l in losses], indent=2))
        return 1 if losses else 0

    if not losses:
        print(f"Nothing dropped: every fact, name, and phrase in {args.original.name} "
              f"still appears somewhere in {args.rebuilt.name}.")
        return 0

    label = {"FACT": "numbers", "NAME": "names", "PHRASE": "phrases"}
    for kind in ("FACT", "NAME", "PHRASE"):
        group = [l for l in losses if l.kind == kind]
        if not group:
            continue
        print(f"\nDropped {label[kind]} ({len(group)}):")
        for l in group:
            where = ", ".join(str(s) for s in l.slides)
            value = l.value if len(l.value) <= 90 else l.value[:87] + "…"
            print(f"  was on slide {where}: {value}")

    print(f"\n{len(losses)} item(s) in {args.original.name} are not in {args.rebuilt.name}.")
    print("Confirm each was cut deliberately. This is a prompt to check, not a verdict:")
    print("reworded content and content moved between slides can both show up here.")
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
