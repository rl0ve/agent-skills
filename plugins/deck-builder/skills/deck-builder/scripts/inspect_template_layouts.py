#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


DEFAULT_TEMPLATE = (
    Path(__file__).resolve().parent.parent / "brand" / "build_base.pptx"
)


def fmt_inches(value) -> str:
    return f"{value / Inches(1):.2f}"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="List PowerPoint layout names, placeholders, and placeholder bounds."
    )
    parser.add_argument(
        "template_path",
        nargs="?",
        default=str(DEFAULT_TEMPLATE),
        help="Path to a .pptx template or deck. Defaults to brand/build_base.pptx in this skill.",
    )
    parser.add_argument(
        "--filter",
        default="",
        help="Only show layouts whose names contain this case-insensitive substring.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    template_path = Path(args.template_path).expanduser().resolve()
    if not template_path.exists():
        raise SystemExit(f"File not found: {template_path}")

    query = args.filter.lower()
    prs = Presentation(str(template_path))

    seen: set[str] = set()
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name in seen:
                continue
            seen.add(layout.name)

            if query and query not in layout.name.lower():
                continue

            print(layout.name)
            for shape in layout.placeholders:
                placeholder_type = getattr(shape.placeholder_format.type, "name", str(shape.placeholder_format.type))
                print(
                    "  - "
                    f"{placeholder_type:<14} | {shape.name} | "
                    f"x={fmt_inches(shape.left)} y={fmt_inches(shape.top)} "
                    f"w={fmt_inches(shape.width)} h={fmt_inches(shape.height)}"
                )
            print()

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
