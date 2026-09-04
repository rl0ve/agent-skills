#!/usr/bin/env python3
from __future__ import annotations

import argparse
import sys
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


BANNED_SNIPPETS = [
    "lorem ipsum",
    "this is not in the master slides",
    "click to add",
    "for bullet",
    "fully blank slide [delete this text]",
    "recommend using a",
]

ALLOWED_EMPTY_PLACEHOLDERS = {
    "DATE",
    "FOOTER",
    "HEADER",
    "SLIDE_NUMBER",
}

SUSPICIOUS_PLACEHOLDERS = {
    "TITLE",
    "CENTER_TITLE",
    "SUBTITLE",
    "BODY",
    "CONTENT",
    "OBJECT",
    "TEXT",
    "PICTURE",
    "MEDIA_CLIP",
    "TABLE",
    "CHART",
    "BITMAP",
    "ORG_CHART",
    "VERTICAL_OBJECT",
}

GENERIC_LAYOUT_NAMES = {
    "Title Slide",
    "Title and Content",
    "Section Header",
    "Two Content",
    "Comparison",
    "Title Only",
    "Blank",
    "Content with Caption",
    "Picture with Caption",
}

BRAND_TEMPLATE_PATH = (
    Path(__file__).resolve().parent.parent / "brand" / "build_base.pptx"
)


def load_template_layout_names(template: Path | None = None) -> set[str]:
    """Layout names from the brand template, or an empty set when it is absent.

    An empty set means the layout-name check cannot run. Callers must say so
    rather than reporting a silent pass.
    """
    path = template or BRAND_TEMPLATE_PATH
    if not path.exists():
        return set()
    prs = Presentation(str(path))
    return {layout.name for master in prs.slide_masters for layout in master.slide_layouts}


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def normalized_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return " ".join((shape.text or "").split())


def check_deck(path: Path, require_brand_template: bool = False,
               brand_template: Path | None = None) -> tuple[list[str], list[str]]:
    prs = Presentation(str(path))
    findings: list[str] = []
    layout_reports: list[str] = []
    known_layouts = load_template_layout_names(brand_template) if require_brand_template else set()
    if require_brand_template and not known_layouts:
        findings.append(
            "Cannot check layout names: no brand template found. Pass --brand-template "
            "or place one at brand/build_base.pptx. Only the stock-layout check ran."
        )

    if len(prs.slide_masters) == 0:
        findings.append("Deck has no slide masters.")

    with ZipFile(path) as archive:
        presentation_xml = archive.read("ppt/presentation.xml").decode("utf-8", errors="ignore")
        if "sectionLst" in presentation_xml:
            findings.append("Deck still contains template section metadata. Strip PowerPoint sections for standalone delivery.")

    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout_name = slide.slide_layout.name
        layout_reports.append(f"Slide {slide_idx}: layout '{layout_name}'")

        if require_brand_template:
            if layout_name in GENERIC_LAYOUT_NAMES:
                findings.append(
                    f"Slide {slide_idx}: layout '{layout_name}' is a stock PowerPoint layout, not one from the brand template."
                )
            elif known_layouts and layout_name not in known_layouts:
                findings.append(
                    f"Slide {slide_idx}: layout '{layout_name}' is not present in the brand template."
                )

        for shape in iter_shapes(slide.shapes):
            text = normalized_text(shape)
            lowered = text.lower()

            for banned in BANNED_SNIPPETS:
                if banned in lowered:
                    findings.append(
                        f"Slide {slide_idx}: found banned helper text '{banned}' in shape '{getattr(shape, 'name', 'unknown')}'."
                    )

            if not shape.is_placeholder:
                continue

            placeholder_type = shape.placeholder_format.type
            placeholder_name = getattr(placeholder_type, "name", str(placeholder_type))

            if text:
                continue

            if placeholder_name in ALLOWED_EMPTY_PLACEHOLDERS:
                continue

            if placeholder_name in SUSPICIOUS_PLACEHOLDERS or "PLACEHOLDER" in shape.name.upper():
                findings.append(
                    (
                        f"Slide {slide_idx}: empty placeholder '{shape.name}' of type '{placeholder_name}' "
                        f"remains in the final slide. Use that placeholder intentionally or switch to a cleaner layout."
                    )
                )

    return findings, layout_reports


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Validate a PowerPoint deck for template cleanliness and brand-layout discipline."
    )
    parser.add_argument("deck_path", help="Path to the deck to validate")
    parser.add_argument(
        "--require-brand-template",
        action="store_true",
        help="Fail if slides do not use layouts present in the brand template.",
    )
    parser.add_argument(
        "--brand-template",
        type=Path,
        help="Path to the brand .pptx to check layout names against. "
             "Defaults to brand/build_base.pptx beside this skill.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    deck_path = Path(args.deck_path).expanduser().resolve()
    if not deck_path.exists():
        print(f"File not found: {deck_path}", file=sys.stderr)
        return 2

    findings, layout_reports = check_deck(
        deck_path,
        require_brand_template=args.require_brand_template,
        brand_template=args.brand_template,
    )
    print("Slide layouts:")
    for layout_report in layout_reports:
        print(f"- {layout_report}")

    if findings:
        print("Template deck validation failed:")
        for finding in findings:
            print(f"- {finding}")
        return 1

    print("Template deck validation passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
